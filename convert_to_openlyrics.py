#!/usr/bin/env python3
"""
PPT/PPTX to OpenLyrics 0.8 Converter

Converts PowerPoint presentation files containing song lyrics to OpenLyrics XML format.
Follows the rules in agent_instructions.md.
"""

import os
from dotenv import load_dotenv
import re
import sys
import zipfile
import shutil
import subprocess
import tempfile
import unicodedata
from pathlib import Path
from datetime import datetime
from xml.sax.saxutils import escape as xml_escape

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Load environment variables from .env file in the script's directory
# override=True ensures .env values take precedence over shell environment
_script_dir = Path(__file__).parent
load_dotenv(_script_dir / '.env', override=True)

# Constants
OPENLYRICS_NS = "http://openlyrics.info/namespace/2009/song"
LIBREOFFICE_PATH = r"C:\Program Files\LibreOffice\program\soffice.exe"
OPENLYRICS_VERSION = "0.8"
AUDIO_EXTENSIONS = {'.mp3', '.wav', '.m4a', '.wma', '.ogg', '.flac'}
SUPPORTED_EXTENSIONS = {'.pptx', '.ppt', '.odp'}


def convert_ppt_to_pptx(ppt_path: Path, temp_dir: Path) -> Path | None:
    """
    Convert legacy .ppt or .odp to .pptx using LibreOffice.
    Returns path to converted file or None on failure.
    """
    if not Path(LIBREOFFICE_PATH).exists():
        return None
    
    try:
        # Run LibreOffice in headless mode to convert to pptx
        result = subprocess.run(
            [
                LIBREOFFICE_PATH,
                '--headless',
                '--convert-to', 'pptx',
                '--outdir', str(temp_dir),
                str(ppt_path)
            ],
            capture_output=True,
            timeout=60
        )
        
        if result.returncode == 0:
            # Find the converted file
            converted = temp_dir / (ppt_path.stem + '.pptx')
            if converted.exists():
                return converted
    except Exception:
        pass
    
    return None

# Patterns for detection
VERSE_MARKER_PATTERN = re.compile(r'^(\d+)[°ª]?\s*strofa\b', re.IGNORECASE)
CHORUS_LABEL_PATTERN = re.compile(r'^(ritornello|rit\.?)\s*:?\s*', re.IGNORECASE)
TITLE_NUMBER_PATTERN = re.compile(r'^(\d+)\s*[-–—.]?\s*')
CHURCH_FILTER_PATTERN = re.compile(r'chiesa.*olgiate', re.IGNORECASE)


class ConversionLog:
    """Maintains a log file for conversion status."""
    
    def __init__(self, log_path: Path):
        self.log_path = log_path
        self._init_log()
    
    def _init_log(self):
        if not self.log_path.exists():
            with open(self.log_path, 'w', encoding='utf-8') as f:
                f.write(f"Conversion Log - Started {datetime.now().isoformat()}\n")
                f.write("=" * 60 + "\n\n")
    
    def log(self, filename: str, status: str, has_media: bool = False, 
            media_extracted: str = None, notes: str = None, changes: list = None):
        with open(self.log_path, 'a', encoding='utf-8') as f:
            f.write(f"File: {filename}\n")
            f.write(f"  Status: {status}\n")
            f.write(f"  Has Media: {'Yes' if has_media else 'No'}\n")
            if media_extracted:
                f.write(f"  Media Extracted: {media_extracted}\n")
            if changes:
                f.write(f"  Changes Applied:\n")
                for change in changes:
                    f.write(f"    - {change}\n")
            if notes:
                f.write(f"  Notes: {notes}\n")
            f.write(f"  Processed: {datetime.now().isoformat()}\n")
            f.write("-" * 40 + "\n")


def clean_text(text: str) -> str:
    """
    Clean text by stripping control chars and normalizing whitespace.
    Also converts all-uppercase lines to sentence case.
    Applies custom text replacements from TEXT_REPLACEMENTS env variable.
    """
    if not text:
        return ""
    
    # Strip ASCII control chars (0x00-0x1F, 0x7F) except newline/tab
    result = []
    for char in text:
        code = ord(char)
        if code == 0x09:  # Tab -> newline
            result.append('\n')
        elif code == 0x0A:  # Newline - keep
            result.append(char)
        elif code == 0x0B or code == 0x0C:  # Vertical tab, form feed -> newline
            result.append('\n')
        elif code == 0x0D:  # Carriage return -> newline
            result.append('\n')
        elif 0x00 <= code <= 0x1F or code == 0x7F:  # Other control chars - remove
            continue
        else:
            result.append(char)
    
    text = ''.join(result)
    
    # Apply custom text replacements from environment variable
    replacements_str = os.getenv('TEXT_REPLACEMENTS', '')
    if replacements_str:
        # Parse replacements: format is "old|new,old2|new2"
        for replacement in replacements_str.split(','):
            replacement = replacement.strip()
            if '|' in replacement:
                old_text, new_text = replacement.split('|', 1)
                text = text.replace(old_text, new_text)
    
    # Normalize multiple newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Trim trailing spaces per line and normalize case
    lines = text.split('\n')
    normalized_lines = []
    for line in lines:
        stripped = line.rstrip()
        if stripped and stripped.isupper():
            # Convert all-uppercase lines to sentence case
            stripped = stripped.capitalize()
        normalized_lines.append(stripped)
    
    return '\n'.join(normalized_lines)


def extract_title_from_filename(filename: str) -> tuple[str, list[str], str | None, str | None]:
    """
    Extract song title from filename, optionally stripping leading numbers and separators.
    Controlled by PPT_REMOVE_NUMBER_PREFIX environment variable.
    If set to 'false', '0', 'no', or 'off' (case insensitive), numbers are kept.
    
    Returns:
        - title: The title in Title Case
        - changes: List of changes made
        - songbook_name: Title without number prefix (for songbook feature)
        - songbook_entry: The number prefix if present (for songbook entry)
    """
    changes = []
    songbook_entry = None
    
    # Remove extension
    name = Path(filename).stem
    original_name = name
    
    # Check for number prefix (always extract for songbook, even if not removing from title)
    number_match = TITLE_NUMBER_PATTERN.match(name)
    if number_match:
        songbook_entry = number_match.group(1)  # Just the number, not separators
    
    # Check environment variable for number prefix removal
    remove_number_prefix = os.getenv('PPT_REMOVE_NUMBER_PREFIX', 'true').lower() not in ('false', '0', 'no', 'off')
    
    # Remove leading numbers and separators (if enabled)
    if remove_number_prefix:
        if number_match:
            removed_prefix = number_match.group(0)
            name = TITLE_NUMBER_PATTERN.sub('', name)
            changes.append(f"Removed number prefix: '{removed_prefix.strip()}'")
    
    # Remove trailing underscores
    if name.endswith('_'):
        name = name.rstrip('_')
        changes.append("Removed trailing underscore(s)")
    
    # Clean up extra spaces
    name = ' '.join(name.split())
    
    # Convert to Title Case
    name = name.strip() or Path(filename).stem
    
    # Custom title case that handles apostrophes and accented chars better
    def title_case_word(word: str) -> str:
        if not word:
            return word
        # Handle words with apostrophes (e.g., "DELL'AMICIZIA" -> "Dell'Amicizia")
        if "'" in word:
            parts = word.split("'")
            return "'".join(part.capitalize() for part in parts)
        return word.capitalize()
    
    words = name.split()
    titled = ' '.join(title_case_word(word) for word in words)
    
    # Songbook name is title without number (always computed)
    name_without_number = TITLE_NUMBER_PATTERN.sub('', original_name).strip()
    name_without_number = ' '.join(name_without_number.split())
    songbook_words = name_without_number.split()
    songbook_name = ' '.join(title_case_word(word) for word in songbook_words)
    
    # Check if title case changed anything
    if titled != name:
        changes.append(f"Converted to Title Case: '{original_name}' → '{titled}'")
    
    return titled, changes, songbook_name, songbook_entry


def extract_text_from_pptx(pptx_path: Path) -> list[list[str]]:
    """
    Extract text from PPTX file, returning list of slides, each slide is a list of text blocks.
    """
    def extract_text_from_shape(shape) -> list[str]:
        """Recursively extract text from a shape, including group shapes."""
        texts = []
        if shape.has_text_frame:
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text
                if para_text.strip():
                    text_parts.append(para_text)
            if text_parts:
                texts.append('\n'.join(text_parts))
        elif hasattr(shape, 'shapes'):  # Group shape
            for sub_shape in shape.shapes:
                texts.extend(extract_text_from_shape(sub_shape))
        return texts
    
    slides_text = []
    
    try:
        prs = Presentation(str(pptx_path))
        
        for slide in prs.slides:
            slide_texts = []
            
            for shape in slide.shapes:
                slide_texts.extend(extract_text_from_shape(shape))
            
            if slide_texts:
                slides_text.append(slide_texts)
    
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from {pptx_path}: {e}")
    
    return slides_text


def extract_media_from_pptx(pptx_path: Path, output_dir: Path, base_name: str) -> str | None:
    """
    Extract embedded audio from PPTX if present.
    Returns the extracted filename or None.
    """
    try:
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            media_files = [n for n in zf.namelist() if n.startswith('ppt/media/')]
            
            for media_file in media_files:
                ext = Path(media_file).suffix.lower()
                if ext in AUDIO_EXTENSIONS:
                    # Check if audio already exists
                    existing_audio = find_existing_audio(output_dir, base_name)
                    if existing_audio:
                        return None  # Skip, already exists
                    
                    # Extract with song title
                    output_name = f"{base_name}{ext}"
                    output_path = output_dir / output_name
                    
                    with zf.open(media_file) as src:
                        with open(output_path, 'wb') as dst:
                            shutil.copyfileobj(src, dst)
                    
                    return output_name
    
    except Exception:
        pass
    
    return None


def find_existing_audio(folder: Path, base_name: str) -> Path | None:
    """
    Check if an audio file matching the base name already exists.
    """
    # Normalize base name for comparison
    base_lower = base_name.lower()
    
    for file in folder.iterdir():
        if file.suffix.lower() in AUDIO_EXTENSIONS:
            file_base, _, _, _ = extract_title_from_filename(file.name)
            if file_base.lower() == base_lower or base_lower in file.stem.lower():
                return file
    
    return None


def is_header_footer(text: str, title: str, all_slides_text: list[list[str]], 
                     filtered_items: list = None) -> bool:
    """
    Check if text is a header/footer that should be ignored.
    Only filters out church names and very short standalone headers.
    Does NOT filter out repeated lyric phrases (like refrains within verses).
    Title filtering is handled separately in parse_slides_to_blocks.
    If filtered_items list is provided, appends the reason for filtering.
    """
    text_stripped = text.strip()
    
    # Church name filter
    if CHURCH_FILTER_PATTERN.search(text_stripped):
        if filtered_items is not None:
            filtered_items.append(f"Removed church name header: '{text_stripped}'")
        return True
    
    # Check if line matches title with number prefix (e.g., "178 - Come il riarso terreno")
    if title:
        text_no_number = TITLE_NUMBER_PATTERN.sub('', text_stripped).strip()
        text_no_number_normalized = _normalize_for_title_comparison(text_no_number)
        title_normalized = _normalize_for_title_comparison(title)
        
        # If after removing number prefix, the line matches the title, it's a header
        if text_no_number_normalized == title_normalized and text_no_number != text_stripped:
            if filtered_items is not None:
                filtered_items.append(f"Removed numbered title header: '{text_stripped}'")
            return True
    
    # Only filter very short text (<=3 words) that appears on most slides as standalone
    # This catches headers like "Page 1" or church names, but not lyric refrains
    words = text_stripped.split()
    if len(words) > 3:
        return False  # Longer text is likely lyrics, not header/footer
    
    # Check if this short text appears on most slides as a standalone element
    occurrence_count = 0
    total_slides = len(all_slides_text)
    
    for slide_texts in all_slides_text:
        for slide_text in slide_texts:
            # Only count if text is standalone (the entire text block) or at very start
            if text_stripped == slide_text.strip():
                occurrence_count += 1
                break
    
    if total_slides > 2 and occurrence_count >= total_slides * 0.7:
        if filtered_items is not None:
            filtered_items.append(f"Removed recurring header/footer: '{text_stripped}' (appeared on {occurrence_count}/{total_slides} slides)")
        return True
    
    return False


def _normalize_for_title_comparison(text: str) -> str:
    """Normalize text for title comparison (unicode, apostrophes, whitespace, punctuation)."""
    import re
    # NFKD decomposition
    normalized = unicodedata.normalize('NFKD', text.lower())
    # Remove combining diacritical marks (accents, etc.)
    normalized = ''.join(c for c in normalized if not unicodedata.combining(c))
    # Remove all types of apostrophes, quotes, and similar characters
    normalized = re.sub(r"[''`'ʼʻˈˊ\u0027\u2019\u2018\u02BC\u02BB\u0060\u00B4]", '', normalized)
    normalized = re.sub(r'[""„‟\u0022\u201C\u201D\u201E\u201F]', '', normalized)
    # Remove common punctuation (commas, periods, exclamation marks, question marks, etc.)
    normalized = re.sub(r'[,;.!?:¡¿]', '', normalized)
    # Normalize whitespace first
    normalized = ' '.join(normalized.split())
    # Normalize "oh" to "o" as standalone word (common in Italian songs)
    normalized = re.sub(r'\boh\b', 'o', normalized)
    return normalized


def _check_if_all_slides_have_title_first(slides_text: list[list[str]], title: str) -> bool:
    """
    Check if ALL slides have the title as the first line (isolated).
    Returns True if all slides start with the title.
    """
    if not slides_text:
        return False
    
    title_normalized = _normalize_for_title_comparison(title)
    title_no_number = _normalize_for_title_comparison(TITLE_NUMBER_PATTERN.sub('', title).strip())
    
    for slide_texts in slides_text:
        if not slide_texts:
            return False  # Empty slide means not all have title
        
        # Get the first text block of the slide
        first_block = slide_texts[0]
        first_block_clean = clean_text(first_block)
        lines = [l.strip() for l in first_block_clean.split('\n') if l.strip()]
        
        if not lines:
            return False  # No lines means not all have title
        
        first_line = lines[0]
        first_line_normalized = _normalize_for_title_comparison(first_line)
        first_line_no_number = _normalize_for_title_comparison(TITLE_NUMBER_PATTERN.sub('', first_line).strip())
        
        # Check if first line matches title (with or without number)
        is_title = (first_line_normalized == title_normalized or 
                   first_line_no_number == title_no_number and first_line_no_number)
        
        if not is_title:
            return False
    
    return True


def _is_title_line(line: str, title: str, is_first_line: bool, is_isolated: bool,
                   all_slides_have_title: bool, slide_index: int, filtered_items: list = None) -> bool:
    """
    Check if a line is a title that should be filtered.
    
    Rules:
    1. Must be the first line of the slide
    2. Must be isolated (not immediately followed by another line)
    3. From slide 2 onwards, only filter if ALL slides have title as first line
    """
    if not is_first_line:
        return False
    
    if not is_isolated:
        return False
    
    # For slides after the first, only filter if all slides have the title
    if slide_index > 0 and not all_slides_have_title:
        return False
    
    title_normalized = _normalize_for_title_comparison(title)
    title_no_number = _normalize_for_title_comparison(TITLE_NUMBER_PATTERN.sub('', title).strip())
    
    line_normalized = _normalize_for_title_comparison(line)
    line_no_number = _normalize_for_title_comparison(TITLE_NUMBER_PATTERN.sub('', line).strip())
    
    # Check if line matches title (with or without number prefix)
    if line_normalized == title_normalized:
        if filtered_items is not None:
            filtered_items.append(f"Removed title from slide {slide_index + 1}: '{line}'")
        return True
    
    if line_no_number == title_no_number and line_no_number:
        if filtered_items is not None:
            filtered_items.append(f"Removed numbered title from slide {slide_index + 1}: '{line}'")
        return True
    
    return False


def parse_slides_to_blocks(slides_text: list[list[str]], title: str, 
                           changes: list = None) -> list[dict]:
    """
    Parse slide text into verse/chorus blocks.
    Returns list of {'type': 'verse'|'chorus', 'lines': [...], 'number': int}
    If changes list is provided, appends information about filtered content.
    """
    blocks = []
    verse_num = 0
    chorus_num = 0
    seen_choruses = {}  # For detecting repeated blocks
    filtered_items = []  # Track what was filtered
    verse_markers_found = []
    chorus_labels_found = []
    
    # Pre-check: do all slides have the title as first isolated line?
    all_slides_have_title = _check_if_all_slides_have_title_first(slides_text, title)
    
    for slide_index, slide_texts in enumerate(slides_text):
        # Combine all text blocks from this slide
        combined_text = '\n'.join(slide_texts)
        combined_text = clean_text(combined_text)
        
        if not combined_text.strip():
            continue
        
        lines = combined_text.split('\n')
        non_empty_lines = [l.strip() for l in lines if l.strip()]
        
        # Special case: Check if first slide contains only the title (possibly split across lines)
        # Filter out headers/footers first to get only content lines
        content_lines = [l for l in non_empty_lines if not is_header_footer(l, title, slides_text)]
        
        if slide_index == 0 and content_lines:
            # Combine all content lines and check if it matches the title
            combined_content = ' '.join(content_lines)
            combined_normalized = _normalize_for_title_comparison(combined_content)
            title_normalized = _normalize_for_title_comparison(title)
            title_no_number = _normalize_for_title_comparison(TITLE_NUMBER_PATTERN.sub('', title).strip())
            
            if combined_normalized == title_normalized or (combined_normalized == title_no_number and title_no_number):
                # First slide is just the title - skip it entirely
                if filtered_items is not None:
                    filtered_items.append(f"Removed title slide: '{combined_content}'")
                continue
        
        filtered_lines = []
        current_is_chorus = False
        explicit_verse_num = None
        is_first_content_line = True  # Track if this is the first content line of the slide
        
        line_position_in_original = 0
        for line_index, line in enumerate(lines):
            line_stripped = line.strip()
            
            if not line_stripped:
                line_position_in_original += 1
                continue
            
            # Check if this line is isolated (first line with a blank line after it, or only line)
            is_isolated = False
            if is_first_content_line:
                if len(non_empty_lines) == 1:
                    # Only one line on the whole slide - it's isolated
                    is_isolated = True
                else:
                    # Check if there's a blank line immediately after this first content line
                    if line_position_in_original + 1 < len(lines) and not lines[line_position_in_original + 1].strip():
                        is_isolated = True
                    # Also check if title is in its own text block
                    elif slide_texts:
                        first_block_clean = clean_text(slide_texts[0])
                        first_block_lines = [l.strip() for l in first_block_clean.split('\n') if l.strip()]
                        if len(first_block_lines) == 1:
                            is_isolated = True
            
            line_position_in_original += 1
            
            # Check for title line (new logic)
            if is_first_content_line and _is_title_line(
                line_stripped, title, is_first_content_line, is_isolated,
                all_slides_have_title, slide_index, filtered_items
            ):
                is_first_content_line = False
                continue
            
            is_first_content_line = False
            
            # Check for header/footer (excluding title - that's handled above)
            if is_header_footer(line_stripped, title, slides_text, filtered_items):
                continue
            
            # Check for verse marker (1° Strofa, 2° Strofa, etc.)
            verse_match = VERSE_MARKER_PATTERN.match(line_stripped)
            if verse_match:
                explicit_verse_num = int(verse_match.group(1))
                if line_stripped not in verse_markers_found:
                    verse_markers_found.append(line_stripped)
                continue  # Skip the marker line
            
            # Check for chorus label
            chorus_match = CHORUS_LABEL_PATTERN.match(line_stripped)
            if chorus_match:
                # If we have accumulated lines, save them as verse first
                if filtered_lines and not current_is_chorus:
                    verse_num += 1
                    actual_num = explicit_verse_num if explicit_verse_num else verse_num
                    blocks.append({
                        'type': 'verse',
                        'lines': filtered_lines.copy(),
                        'number': actual_num
                    })
                    filtered_lines = []
                    explicit_verse_num = None
                
                current_is_chorus = True
                # Track the chorus label found
                label_found = chorus_match.group(0).strip()
                if label_found not in chorus_labels_found:
                    chorus_labels_found.append(label_found)
                # Get remaining text after chorus label
                remaining = CHORUS_LABEL_PATTERN.sub('', line_stripped).strip()
                if remaining:
                    filtered_lines.append(remaining)
                continue
            
            filtered_lines.append(line_stripped)
        
        if not filtered_lines:
            continue
        
        # Check for repeated block (auto-detect chorus)
        block_key = '\n'.join(filtered_lines).lower()
        
        if current_is_chorus:
            if block_key not in seen_choruses:
                chorus_num += 1
                seen_choruses[block_key] = chorus_num
            
            blocks.append({
                'type': 'chorus',
                'lines': filtered_lines,
                'number': seen_choruses[block_key]
            })
        elif block_key in seen_choruses:
            # This is a repeated block - it's a chorus
            blocks.append({
                'type': 'chorus',
                'lines': filtered_lines,
                'number': seen_choruses[block_key]
            })
        else:
            # Check if this block repeats later (auto-detect chorus)
            is_repeated = False
            for future_slide_texts in slides_text[slides_text.index(slide_texts) + 1:]:
                future_combined = clean_text('\n'.join(future_slide_texts))
                future_lines = [l.strip() for l in future_combined.split('\n') if l.strip()]
                # Filter out headers
                future_lines = [l for l in future_lines if not is_header_footer(l, title, slides_text)]
                future_key = '\n'.join(future_lines).lower()
                
                if future_key == block_key:
                    is_repeated = True
                    break
            
            if is_repeated and len(filtered_lines) <= 8:  # Choruses are usually short
                if block_key not in seen_choruses:
                    chorus_num += 1
                    seen_choruses[block_key] = chorus_num
                
                blocks.append({
                    'type': 'chorus',
                    'lines': filtered_lines,
                    'number': seen_choruses[block_key]
                })
            else:
                verse_num += 1
                actual_num = explicit_verse_num if explicit_verse_num else verse_num
                blocks.append({
                    'type': 'verse',
                    'lines': filtered_lines,
                    'number': actual_num
                })
    
    # Add tracked changes to the changes list
    if changes is not None:
        # Remove duplicates from filtered_items
        seen = set()
        for item in filtered_items:
            if item not in seen:
                changes.append(item)
                seen.add(item)
        
        if verse_markers_found:
            changes.append(f"Detected verse markers: {', '.join(verse_markers_found)}")
        
        if chorus_labels_found:
            changes.append(f"Detected chorus labels: {', '.join(chorus_labels_found)}")
        
        if seen_choruses:
            changes.append(f"Auto-detected {len(seen_choruses)} chorus pattern(s) from repeated blocks")
    
    return blocks


def deduplicate_blocks(blocks: list[dict]) -> list[dict]:
    """
    Remove consecutive duplicate blocks (same content).
    """
    if not blocks:
        return blocks
    
    result = [blocks[0]]
    
    for block in blocks[1:]:
        prev_key = '\n'.join(result[-1]['lines']).lower()
        curr_key = '\n'.join(block['lines']).lower()
        
        if curr_key != prev_key:
            result.append(block)
    
    return result


def generate_openlyrics_xml(title: str, blocks: list[dict], 
                           songbook_name: str = None, songbook_entry: str = None) -> str:
    """
    Generate OpenLyrics 0.8 XML from parsed blocks.
    
    Args:
        title: Song title
        blocks: List of verse/chorus blocks
        songbook_name: Optional songbook name (typically title without number)
        songbook_entry: Optional entry number in the songbook
    """
    # Deduplicate consecutive identical blocks
    blocks = deduplicate_blocks(blocks)
    
    # Renumber verses and choruses sequentially
    verse_counter = 0
    chorus_counter = 0
    seen_verses = {}
    seen_choruses = {}
    
    verses_xml = []
    
    for block in blocks:
        lines_escaped = [xml_escape(line) for line in block['lines']]
        lines_content = '<br />'.join(lines_escaped)
        
        if block['type'] == 'chorus':
            key = '\n'.join(block['lines']).lower()
            if key not in seen_choruses:
                chorus_counter += 1
                seen_choruses[key] = chorus_counter
            name = f"c{seen_choruses[key]}"
        else:
            key = '\n'.join(block['lines']).lower()
            if key not in seen_verses:
                verse_counter += 1
                seen_verses[key] = verse_counter
            name = f"v{seen_verses[key]}"
        
        verses_xml.append(f'    <verse name="{name}">\n      <lines>{lines_content}</lines>\n    </verse>')
    
    title_escaped = xml_escape(title)
    
    # Build songbooks section if enabled
    songbooks_xml = ""
    if songbook_name:
        songbook_name_escaped = xml_escape(songbook_name)
        if songbook_entry:
            songbooks_xml = f'''\n    <songbooks>
      <songbook name="{songbook_name_escaped}" entry="{xml_escape(songbook_entry)}"/>
    </songbooks>'''
        else:
            songbooks_xml = f'''\n    <songbooks>
      <songbook name="{songbook_name_escaped}"/>
    </songbooks>'''
    
    xml = f'''<?xml version="1.0" encoding="UTF-8"?>
<song version="{OPENLYRICS_VERSION}" xmlns="{OPENLYRICS_NS}">
  <properties>
    <titles>
      <title>{title_escaped}</title>
    </titles>{songbooks_xml}
  </properties>
  <lyrics>
{chr(10).join(verses_xml)}
  </lyrics>
</song>'''
    
    return xml


def convert_file(input_path: Path, output_dir: Path, log: ConversionLog) -> bool:
    """
    Convert a single PPT/PPTX file to OpenLyrics XML.
    Returns True if successful.
    """
    filename = input_path.name
    title, title_changes, songbook_name, songbook_entry = extract_title_from_filename(filename)
    changes = title_changes.copy()  # Track all changes
    temp_pptx = None
    
    # Check if songbook feature is enabled
    add_songbook = os.getenv('ADD_SONG_BOOK', 'false').lower() in ('true', '1', 'yes', 'on')
    
    try:
        # Extract text - handle different formats
        openlyrics_dir = output_dir / "openlyrics"
        openlyrics_dir.mkdir(exist_ok=True)
        raw_text_dir = output_dir / "raw text"
        raw_text_dir.mkdir(exist_ok=True)
        
        if input_path.suffix.lower() == '.pptx':
            slides_text = extract_text_from_pptx(input_path)
            actual_pptx = input_path
            
            # Write raw extracted text to txt file for debugging
            raw_txt_path = raw_text_dir / f"{title}.txt"
            with open(raw_txt_path, 'w', encoding='utf-8') as f:
                for i, slide in enumerate(slides_text, 1):
                    f.write(f"Slide {i}:\n")
                    for block in slide:
                        # Replace vertical tabs with newlines for better readability
                        clean_block = block.replace('\v', '\n')
                        f.write(clean_block + "\n")
                    f.write("\n")
        elif input_path.suffix.lower() in {'.ppt', '.odp'}:
            # Legacy PPT or ODP - convert using LibreOffice first
            changes.append(f"Converted from legacy {input_path.suffix} format via LibreOffice")
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_dir_path = Path(temp_dir)
                converted_path = convert_ppt_to_pptx(input_path, temp_dir_path)
                
                if converted_path and converted_path.exists():
                    slides_text = extract_text_from_pptx(converted_path)
                    actual_pptx = converted_path
                else:
                    # Fallback: try direct extraction (rarely works for .ppt)
                    try:
                        slides_text = extract_text_from_pptx(input_path)
                        actual_pptx = input_path
                    except Exception:
                        log.log(filename, "SKIPPED", notes="LibreOffice conversion failed and direct extraction not possible")
                        return False
                
                # Write raw extracted text to txt file for debugging
                raw_text_dir = output_dir / "raw text"
                raw_text_dir.mkdir(exist_ok=True)
                raw_txt_path = raw_text_dir / f"{title}.txt"
                with open(raw_txt_path, 'w', encoding='utf-8') as f:
                    for i, slide in enumerate(slides_text, 1):
                        f.write(f"Slide {i}:\n")
                        for block in slide:
                            # Replace vertical tabs with newlines for better readability
                            clean_block = block.replace('\v', '\n')
                            f.write(clean_block + "\n")
                        f.write("\n")
                
                if not slides_text:
                    log.log(filename, "SKIPPED", notes="No text content found after conversion")
                    return False
                
                # Parse into blocks
                blocks = parse_slides_to_blocks(slides_text, title, changes)
                
                if not blocks:
                    log.log(filename, "SKIPPED", notes="No verse/chorus blocks detected")
                    return False
                
                # Generate XML
                xml_content = generate_openlyrics_xml(
                    title, blocks,
                    songbook_name=songbook_name if add_songbook else None,
                    songbook_entry=songbook_entry if add_songbook else None
                )
                
                # Write XML file
                openlyrics_dir = output_dir / "openlyrics"
                openlyrics_dir.mkdir(exist_ok=True)
                output_path = openlyrics_dir / f"{title}.xml"
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(xml_content)
                
                # Extract media to openlyrics folder
                media_extracted = None
                has_media = False
                if converted_path and converted_path.exists():
                    media_extracted = extract_media_from_pptx(converted_path, openlyrics_dir, title)
                    has_media = media_extracted is not None
                
                log.log(filename, "SUCCESS", has_media=has_media, media_extracted=media_extracted,
                       changes=changes if changes else None)
                return True
        else:
            log.log(filename, "SKIPPED", notes=f"Unsupported format: {input_path.suffix}")
            return False
        
        if not slides_text:
            log.log(filename, "SKIPPED", notes="No text content found")
            return False
        
        # Parse into blocks
        blocks = parse_slides_to_blocks(slides_text, title, changes)
        
        if not blocks:
            log.log(filename, "SKIPPED", notes="No verse/chorus blocks detected")
            return False
        
        # Generate XML
        xml_content = generate_openlyrics_xml(
            title, blocks,
            songbook_name=songbook_name if add_songbook else None,
            songbook_entry=songbook_entry if add_songbook else None
        )
        
        # Write XML file
        output_path = openlyrics_dir / f"{title}.xml"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(xml_content)
        
        # Extract media if present
        media_extracted = None
        has_media = False
        
        if input_path.suffix.lower() == '.pptx':
            media_extracted = extract_media_from_pptx(input_path, openlyrics_dir, title)
            has_media = media_extracted is not None
        
        log.log(filename, "SUCCESS", has_media=has_media, media_extracted=media_extracted,
               changes=changes if changes else None)
        return True
    
    except Exception as e:
        log.log(filename, "ERROR", notes=str(e))
        return False


def main():
    """Main entry point."""
    input_base = Path("input")
    output_base = Path("output")
    input_done = Path("input_done")
    
    if not input_base.exists():
        print("Error: input folder not found")
        sys.exit(1)
    
    output_base.mkdir(exist_ok=True)
    input_done.mkdir(exist_ok=True)
    
    total_success = 0
    total_skip = 0
    
    for subfolder in sorted(input_base.iterdir()):
        if not subfolder.is_dir():
            continue
        
        output_folder = output_base / subfolder.name
        output_folder.mkdir(exist_ok=True)
        
        log = ConversionLog(output_folder / "conversion_log.txt")
        
        ppt_files = []
        for ext in SUPPORTED_EXTENSIONS:
            ppt_files.extend(subfolder.glob(f"*{ext}"))
        
        ppt_files.sort(key=lambda p: p.name.lower())
        
        print(f"Processing folder: {subfolder.name}")
        print(f"Found {len(ppt_files)} presentation files")
        print("-" * 40)
        
        success_count = 0
        skip_count = 0
        
        for ppt_file in ppt_files:
            print(f"  Processing: {ppt_file.name}")
            
            if convert_file(ppt_file, output_folder, log):
                success_count += 1
                print("    [OK] Converted successfully")
            else:
                skip_count += 1
                print("    [SKIP] Skipped or failed")
        
        print(f"  Folder {subfolder.name}: {success_count} success, {skip_count} skipped")
        print()
        
        # Move processed folder to input_done
        try:
            shutil.move(str(subfolder), str(input_done / subfolder.name))
            print(f"  Moved folder '{subfolder.name}' to input_done")
        except Exception as e:
            print(f"  Warning: Could not move folder '{subfolder.name}' to input_done: {e}")
        
        total_success += success_count
        total_skip += skip_count
    
    # Also process any files directly in input_base
    ppt_files = []
    for ext in SUPPORTED_EXTENSIONS:
        ppt_files.extend(f for f in input_base.iterdir() if f.is_file() and f.suffix.lower() == ext)
    
    if ppt_files:
        output_folder = output_base / "root_files"
        output_folder.mkdir(exist_ok=True)
        
        log = ConversionLog(output_folder / "conversion_log.txt")
        
        ppt_files.sort(key=lambda p: p.name.lower())
        
        print(f"Processing root files")
        print(f"Found {len(ppt_files)} presentation files")
        print("-" * 40)
        
        success_count = 0
        skip_count = 0
        
        for ppt_file in ppt_files:
            print(f"  Processing: {ppt_file.name}")
            
            if convert_file(ppt_file, output_folder, log):
                success_count += 1
                print("    [OK] Converted successfully")
                # Move processed file to input_done
                try:
                    shutil.move(str(ppt_file), str(input_done / ppt_file.name))
                    print(f"    Moved file '{ppt_file.name}' to input_done")
                except Exception as e:
                    print(f"    Warning: Could not move file '{ppt_file.name}' to input_done: {e}")
            else:
                skip_count += 1
                print("    [SKIP] Skipped or failed")
        
        print(f"  Root files: {success_count} success, {skip_count} skipped")
        print()
        
        total_success += success_count
        total_skip += skip_count
    
    print("-" * 50)
    print("Conversion complete!")
    print(f"  Total Success: {total_success}")
    print(f"  Total Skipped/Errors: {total_skip}")


if __name__ == "__main__":
    main()
